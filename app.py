import streamlit as st
import pandas as pd
from pptx import Presentation
import os
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from email.mime.text import MIMEText

# Function to modify the PPTX certificate template and save it as a new PPTX file
def modify_pptx(template_path, name, output_pptx_path):
    prs = Presentation(template_path)
    
    # Replace {{Name}} placeholder in each slide of the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                shape.text = shape.text.replace('{{Name}}', name)
    
    # Save modified PPTX file
    prs.save(output_pptx_path)

# Function to convert PPTX to PDF using LibreOffice (or another method)
def convert_pptx_to_pdf(input_pptx_path):
    output_pdf_path = input_pptx_path.replace('.pptx', '.pdf')
    command = f'libreoffice --headless --convert-to pdf "{input_pptx_path}" --outdir "{os.path.dirname(input_pptx_path)}"'
    os.system(command)
    return output_pdf_path

# Function to send email with attachment
def send_email(recipient_email, subject, body, attachment_path):
    sender_email = "barathvikraman.projects@gmail.com"  # Your email address
    sender_password = "stlz pqqt prst nfha"       # Your email password or app password

    # Create a multipart message and set headers
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = recipient_email
    msg['Subject'] = subject

    # Attach body text to email
    msg.attach(MIMEText(body.format(name=recipient_email.split('@')[0]), 'plain'))

    # Attach the certificate file (PDF)
    try:
        with open(attachment_path, "rb") as attachment:
            part = MIMEBase("application", "octet-stream")
            part.set_payload(attachment.read())
            encoders.encode_base64(part)
            part.add_header(
                "Content-Disposition",
                f"attachment; filename= {os.path.basename(attachment_path)}",
            )
            msg.attach(part)

        # Send email via SMTP server
        with smtplib.SMTP('smtp.gmail.com', 587) as server:
            server.starttls()  # Upgrade the connection to secure
            server.login(sender_email, sender_password)
            server.send_message(msg)
            st.success(f"Email sent to {recipient_email}")
    
    except FileNotFoundError:
        st.error(f"Attachment file not found: {attachment_path}")
    except Exception as e:
        st.error(f"Failed to send email to {recipient_email}: {str(e)}")

def main():
    st.title("Certificate Generator")

    # File upload section for PPTX and Excel files
    template_file = st.file_uploader("Upload PPTX Template", type=["pptx"])
    data_file = st.file_uploader("Upload Excel File", type=["xlsx"])

    email_subject = st.text_input("Email Subject")
    email_body = st.text_area("Email Body")

    if st.button("Generate Certificates"):
        if template_file and data_file:
            # Create a unique folder based on task name or timestamp
            task_name = st.text_input("Enter Task Name", value="Certificates")
            folder_name = f'uploads/{task_name}'
            os.makedirs(folder_name, exist_ok=True)

            # Save uploaded files temporarily in uploads directory
            template_path = os.path.join(folder_name, template_file.name)
            data_path = os.path.join(folder_name, data_file.name)

            with open(template_path, "wb") as f:
                f.write(template_file.getbuffer())
            with open(data_path, "wb") as f:
                f.write(data_file.getbuffer())

            # Read Excel data containing names and emails
            df = pd.read_excel(data_path)

            for index, row in df.iterrows():
                name = row['Name']  # Adjust based on your Excel column name
                
                output_pptx_filename = f'certificate_{index + 1}.pptx'
                output_pptx_filepath = os.path.join(folder_name, output_pptx_filename)
                
                # Modify the PPTX template and save it with the name replaced
                modify_pptx(template_path, name, output_pptx_filepath)

                # Convert modified PPTX to PDF (make sure LibreOffice is installed)
                pdf_filepath = convert_pptx_to_pdf(output_pptx_filepath)

                # Check if PDF was created successfully before sending email
                if os.path.exists(pdf_filepath):
                    # Send email with attachment (assuming there's an Email column in Excel)
                    recipient_email = row['Email']  # Adjust based on your Excel column name
                    send_email(recipient_email, email_subject, email_body, pdf_filepath)
                else:
                    st.error(f"Failed to create PDF for {name}")

            st.success("All certificates generated and emailed successfully!")

if __name__ == "__main__":
    main()
